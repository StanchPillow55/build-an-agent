"""
Slide generator module for converting curriculum plans to PowerPoint presentations.

Converts a curriculum JSON (from curriculum_planner.plan_curriculum) into a .pptx deck
with title slide, content slides with images, and assessment slide.
"""

import tempfile
from pathlib import Path
from typing import Dict, Any, Optional

import requests
from PIL import Image
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor


def sanitize_keyword(text: str) -> str:
    """Sanitize text for use as an image search keyword."""
    # Remove special characters and convert to lowercase
    sanitized = "".join(c for c in text if c.isalnum() or c.isspace())
    # Replace multiple spaces with single space and strip
    sanitized = " ".join(sanitized.split())
    # Convert to lowercase and replace spaces with +
    return sanitized.lower().replace(" ", "+")


def download_image(keyword: str, max_width: int = 800) -> Optional[Path]:
    """
    Download an image from Unsplash based on keyword and resize it.

    Args:
        keyword: Search keyword for the image
        max_width: Maximum width for the resized image

    Returns:
        Path to the downloaded and resized image, or None if failed
    """
    try:
        # Construct Unsplash URL with multiple fallbacks
        sanitized_keyword = sanitize_keyword(keyword)
        urls = [
            f"https://source.unsplash.com/800x600/?{sanitized_keyword}",
            f"https://source.unsplash.com/featured/?{sanitized_keyword}",
            "https://source.unsplash.com/800x600/?education",  # Fallback
        ]

        response = None
        for url in urls:
            try:
                response = requests.get(
                    url,
                    timeout=10,
                    stream=True,
                    headers={
                        "User-Agent": "Mozilla/5.0 (Windows NT 10.0; Win64; x64) AppleWebKit/537.36"
                    },
                )
                response.raise_for_status()
                break
            except requests.exceptions.RequestException:
                continue

        if response is None:
            return None

        # Create temporary file
        temp_file = tempfile.NamedTemporaryFile(delete=False, suffix=".jpg")

        # Write image data
        for chunk in response.iter_content(chunk_size=8192):
            temp_file.write(chunk)
        temp_file.close()

        # Resize image while maintaining aspect ratio
        with Image.open(temp_file.name) as img:
            if img.width > max_width:
                # Calculate new height maintaining aspect ratio
                ratio = max_width / img.width
                new_height = int(img.height * ratio)
                img_resized = img.resize(
                    (max_width, new_height), Image.Resampling.LANCZOS
                )

                # Save resized image
                resized_path = Path(temp_file.name).with_suffix(".resized.jpg")
                img_resized.save(resized_path, "JPEG", quality=85)

                # Remove original temp file
                Path(temp_file.name).unlink()

                return resized_path
            else:
                return Path(temp_file.name)

    except Exception as e:
        print(f"Warning: Failed to download image for '{keyword}': {e}")
        return None


def create_title_slide(
    prs: Presentation, lesson_title: str, learning_objectives: list
) -> None:
    """Create the title slide with lesson title and learning objectives."""
    slide_layout = prs.slide_layouts[0]  # Title slide layout
    slide = prs.slides.add_slide(slide_layout)

    # Set title
    title = slide.shapes.title
    title.text = lesson_title

    # Set subtitle with objectives
    subtitle = slide.placeholders[1]
    subtitle.text = "Learning Objectives:"

    # Add objectives as bullet points
    objectives_text = subtitle.text + "\n\n"
    for i, obj in enumerate(learning_objectives, 1):
        objectives_text += f"• {obj}\n"

    subtitle.text = objectives_text.strip()

    # Format subtitle
    for paragraph in subtitle.text_frame.paragraphs:
        paragraph.font.size = Pt(18)
        paragraph.font.color.rgb = RGBColor(64, 64, 64)


def create_content_slide(
    prs: Presentation, title: str, description: str, image_path: Optional[Path] = None
) -> None:
    """Create a content slide with title, description, and optional image."""
    slide_layout = prs.slide_layouts[1]  # Title and content layout
    slide = prs.slides.add_slide(slide_layout)

    # Set title
    title_shape = slide.shapes.title
    title_shape.text = title

    # Set content
    content = slide.placeholders[1]
    content.text = description

    # Format content
    for paragraph in content.text_frame.paragraphs:
        paragraph.font.size = Pt(20)
        paragraph.font.color.rgb = RGBColor(64, 64, 64)

    # Add image if available
    if image_path and image_path.exists():
        try:
            # Position image on the right side
            left = Inches(6)
            top = Inches(2)
            width = Inches(3.5)

            slide.shapes.add_picture(str(image_path), left, top, width=width)

            # Adjust content placeholder to make room for image
            content.left = Inches(0.5)
            content.width = Inches(5)

        except Exception as e:
            print(f"Warning: Failed to add image to slide: {e}")
        finally:
            # Clean up temporary image file
            try:
                image_path.unlink()
            except Exception:
                pass


def create_assessment_slide(prs: Presentation, assessments: list) -> None:
    """Create the final assessment slide."""
    slide_layout = prs.slide_layouts[1]  # Title and content layout
    slide = prs.slides.add_slide(slide_layout)

    # Set title
    title = slide.shapes.title
    title.text = "Assessments"

    # Set content with assessments
    content = slide.placeholders[1]
    assessment_text = "Suggested Assessment Methods:\n\n"

    for i, assessment in enumerate(assessments, 1):
        assessment_text += f"• {assessment}\n"

    content.text = assessment_text.strip()

    # Format content
    for paragraph in content.text_frame.paragraphs:
        paragraph.font.size = Pt(20)
        paragraph.font.color.rgb = RGBColor(64, 64, 64)


def create_deck(plan: Dict[str, Any], output_path: Path) -> None:
    """
    Convert a curriculum plan JSON to a PowerPoint presentation.

    Args:
        plan: Curriculum plan dictionary from curriculum_planner.plan_curriculum
        output_path: Path where the .pptx file should be saved
    """
    # Create presentation
    prs = Presentation()

    # Extract data from plan
    lesson_title = plan.get("lesson_title", "Lesson Plan")
    learning_objectives = plan.get("learning_objectives", [])
    content_outline = plan.get("content_outline", [])
    suggested_assessments = plan.get("suggested_assessments", [])

    # Create title slide
    create_title_slide(prs, lesson_title, learning_objectives)

    # Create content slides
    for content_item in content_outline:
        title = content_item.get("title", "Content Section")
        description = content_item.get("description", "")

        # Download image based on title
        image_path = download_image(title)

        # Create slide
        create_content_slide(prs, title, description, image_path)

    # Create assessment slide
    if suggested_assessments:
        create_assessment_slide(prs, suggested_assessments)

    # Save presentation
    output_path.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(output_path))

    print(f"✅ PowerPoint deck saved to {output_path}")


def main():
    """Test the slide generator with sample data."""
    # Sample curriculum plan for testing
    sample_plan = {
        "lesson_title": "Introduction to Environmental Science",
        "learning_objectives": [
            "Students will define what an ecosystem is",
            "Students will identify biotic and abiotic factors",
            "Students will explain food chains and energy flow",
        ],
        "content_outline": [
            {
                "title": "What is an Ecosystem?",
                "description": "Introduce the concept of ecosystems using local examples and interactive discussions.",
            },
            {
                "title": "Living vs Non-Living Components",
                "description": "Explore biotic and abiotic factors through hands-on activities and observations.",
            },
            {
                "title": "Energy Flow in Ecosystems",
                "description": "Demonstrate food chains and energy transfer concepts using visual models.",
            },
        ],
        "suggested_assessments": [
            "Ecosystem components identification worksheet",
            "Food chain construction activity",
            "Exit ticket with key vocabulary terms",
        ],
    }

    # Create test presentation
    output_path = Path("test_presentation.pptx")
    create_deck(sample_plan, output_path)
    print(f"Test presentation created: {output_path}")


if __name__ == "__main__":
    main()
